from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

# Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "FMT_project 실행 트러블슈팅 및 해결 과정"
subtitle.text = "앱 실행 오류 원인 분석 및 해결 방안"

# Slide 1: Python Version Issue
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "1. 가상 환경 Python 버전 불일치 문제"
tf = body_shape.text_frame
tf.text = "증상: .venv\\Scripts\\python을 찾지 못하거나 버전 에러 발생"
p = tf.add_paragraph()
p.text = "원인: 시스템에 설치된 파이썬(Python 3.12)과 가상 환경(Python 3.10)의 경로 불일치"
p.level = 1
p = tf.add_paragraph()
p.text = "해결: 기존 .venv 삭제 후 Python 3.12로 새로운 가상 환경 생성 및 패키지 재설치"
p.level = 1

# Slide 2: Smart App Control
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "2. 윈도우 스마트 앱 컨트롤 차단 문제"
tf = body_shape.text_frame
tf.text = "증상: matplotlib, pandas 등 라이브러리 로드 시 'DLL load failed' 에러 발생"
p = tf.add_paragraph()
p.text = "원인: Windows 11 스마트 앱 컨트롤이 서명되지 않은 라이브러리(.pyd) 실행을 악성코드로 오인하여 차단"
p.level = 1
p = tf.add_paragraph()
p.text = "해결 방안:"
p.level = 1
p = tf.add_paragraph()
p.text = "1. 윈도우 설정에서 '스마트 앱 컨트롤' 끄기"
p.level = 2
p = tf.add_paragraph()
p.text = "2. PC 재부팅 후 파워쉘에서 앱 재실행"
p.level = 2

# Slide 3: Correct Execution Command
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "3. 올바른 앱 실행 명령어 및 경로"
tf = body_shape.text_frame
tf.text = "증상: 파워쉘에서 실행 시 CommandNotFoundException 발생"
p = tf.add_paragraph()
p.text = "원인: 실행 위치가 프로젝트 폴더가 아닌 C:\\WINDOWS\\system32로 설정됨"
p.level = 1
p = tf.add_paragraph()
p.text = "해결: cd 명령어로 프로젝트 폴더로 이동 후 실행"
p.level = 1
p = tf.add_paragraph()
p.text = "최종 실행 명령어 (한 줄 요약):"
p.level = 1
p = tf.add_paragraph()
p.text = "cd C:\\Users\\handi\\Desktop\\pages\\FMT_project ; .\\.venv\\Scripts\\python.exe -m streamlit run app.py --server.port 8501"
p.level = 2
p.font.size = Pt(14)

prs.save("FMT_Troubleshooting.pptx")
